import json
import boto3
import os
import time
import datetime
from io import BytesIO
import PyPDF2
import csv
import traceback
import re
from urllib import parse

from botocore.config import Config
from multiprocessing import Process, Pipe
from googleapiclient.discovery import build

from langchain.chains import LLMChain
from langchain.prompts import PromptTemplate
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.chains.summarize import load_summarize_chain
from langchain.chains import ConversationChain
from langchain.memory import ConversationBufferWindowMemory
from langchain.callbacks.streaming_stdout import StreamingStdOutCallbackHandler

from langchain_community.vectorstores.faiss import FAISS
from langchain_community.vectorstores.opensearch_vector_search import OpenSearchVectorSearch
from langchain_community.embeddings import BedrockEmbeddings
from langchain_community.docstore.document import Document

from langchain_core.prompts import MessagesPlaceholder, ChatPromptTemplate
from langchain_core.messages import HumanMessage, SystemMessage
from langchain_aws import ChatBedrock

s3 = boto3.client('s3')
s3_bucket = os.environ.get('s3_bucket') # bucket name
s3_prefix = os.environ.get('s3_prefix')
meta_prefix = "metadata/"
callLogTableName = os.environ.get('callLogTableName')
profile_of_LLMs = json.loads(os.environ.get('profile_of_LLMs'))
isReady = False   

opensearch_account = os.environ.get('opensearch_account')
opensearch_passwd = os.environ.get('opensearch_passwd')
enableReference = os.environ.get('enableReference', 'false')
opensearch_url = os.environ.get('opensearch_url')
path = os.environ.get('path')
doc_prefix = s3_prefix+'/'

roleArn = os.environ.get('roleArn')
top_k = int(os.environ.get('numberOfRelevantDocs', '8'))
selected_LLM = 0
MSG_LENGTH = 100
MSG_HISTORY_LENGTH = 20
speech_generation = True
history_length = 0
token_counter_history = 0
allowDualSearching = os.environ.get('allowDualSearching')
allowInternetSearching = os.environ.get('allowInternetSearching')
 
# google search api
googleApiSecret = os.environ.get('googleApiSecret')
secretsmanager = boto3.client('secretsmanager')
try:
    get_secret_value_response = secretsmanager.get_secret_value(
        SecretId=googleApiSecret
    )
    #print('get_secret_value_response: ', get_secret_value_response)
    secret = json.loads(get_secret_value_response['SecretString'])
    #print('secret: ', secret)
    google_api_key = secret['google_api_key']
    google_cse_id = secret['google_cse_id']
    #print('google_cse_id: ', google_cse_id)    

except Exception as e:
    raise e

# websocket
connection_url = os.environ.get('connection_url')
client = boto3.client('apigatewaymanagementapi', endpoint_url=connection_url)
print('connection_url: ', connection_url)

HUMAN_PROMPT = "\n\nHuman:"
AI_PROMPT = "\n\nAssistant:"

map_chain = dict() 

# Multi-LLM
def get_chat(profile_of_LLMs, selected_LLM):
    profile = profile_of_LLMs[selected_LLM]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    print(f'LLM: {selected_LLM}, bedrock_region: {bedrock_region}, modelId: {modelId}')
    maxOutputTokens = int(profile['maxOutputTokens'])
                          
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }            
        )
    )
    parameters = {
        "max_tokens":maxOutputTokens,     
        "temperature":0.1,
        "top_k":250,
        "top_p":0.9,
        "stop_sequences": [HUMAN_PROMPT]
    }
    # print('parameters: ', parameters)

    chat = ChatBedrock(   
        model_id=modelId,
        client=boto3_bedrock, 
        model_kwargs=parameters,
    )       
    
    return chat

def get_embedding(profile_of_LLMs, selected_LLM):
    profile = profile_of_LLMs[selected_LLM]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    print(f'Embedding: {selected_LLM}, bedrock_region: {bedrock_region}, modelId: {modelId}')
    
    # bedrock   
    boto3_bedrock = boto3.client(
        service_name='bedrock-runtime',
        region_name=bedrock_region,
        config=Config(
            retries = {
                'max_attempts': 30
            }            
        )
    )
    
    bedrock_embedding = BedrockEmbeddings(
        client=boto3_bedrock,
        region_name = bedrock_region,
        model_id = 'amazon.titan-embed-text-v1' 
    )  
    
    return bedrock_embedding

def sendMessage(id, body):
    try:
        client.post_to_connection(
            ConnectionId=id, 
            Data=json.dumps(body)
        )
    except Exception:
        err_msg = traceback.format_exc()
        print('err_msg: ', err_msg)
        raise Exception ("Not able to send a message")

def sendResultMessage(connectionId, requestId, msg):    
    result = {
        'request_id': requestId,
        'msg': msg,
        'status': 'completed'
    }
    #print('debug: ', json.dumps(debugMsg))
    sendMessage(connectionId, result)

def sendDebugMessage(connectionId, requestId, msg):
    debugMsg = {
        'request_id': requestId,
        'msg': msg,
        'status': 'debug'
    }
    #print('debug: ', json.dumps(debugMsg))
    sendMessage(connectionId, debugMsg)

def sendErrorMessage(connectionId, requestId, msg):
    errorMsg = {
        'request_id': requestId,
        'msg': msg,
        'status': 'error'
    }
    print('error: ', json.dumps(errorMsg))
    sendMessage(connectionId, errorMsg)

def isKorean(text):
    # check korean
    pattern_hangul = re.compile('[\u3131-\u3163\uac00-\ud7a3]+')
    word_kor = pattern_hangul.search(str(text))
    # print('word_kor: ', word_kor)

    if word_kor and word_kor != 'None':
        print('Korean: ', word_kor)
        return True
    else:
        print('Not Korean: ', word_kor)
        return False

def general_conversation(connectionId, requestId, chat, query):
    if isKorean(query)==True :
        system = (
            "다음의 Human과 Assistant의 친근한 이전 대화입니다. Assistant은 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다."
        )
    else: 
        system = (
            "Using the following conversation, answer friendly for the newest question. If you don't know the answer, just say that you don't know, don't try to make up an answer. You will be acting as a thoughtful advisor."
        )
    
    human = "{input}"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), MessagesPlaceholder(variable_name="history"), ("human", human)])
    print('prompt: ', prompt)
    
    history = memory_chain.load_memory_variables({})["chat_history"]
    print('memory_chain: ', history)
                
    chain = prompt | chat    
    try: 
        isTyping(connectionId, requestId)  
        stream = chain.invoke(
            {
                "history": history,
                "input": query,
            }
        )
        msg = readStreamMsg(connectionId, requestId, stream.content)    
                            
        msg = stream.content
        print('msg: ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
            
        sendErrorMessage(connectionId, requestId, err_msg)    
        raise Exception ("Not able to request to LLM")
    
    return msg
    
def store_document_for_faiss(docs, vectorstore_faiss):
    print('store document into faiss')    
    vectorstore_faiss.add_documents(docs)       
    print('uploaded into faiss')

from opensearchpy import OpenSearch
def delete_index_if_exist(index_name):
    client = OpenSearch(
        hosts = [{
            'host': opensearch_url.replace("https://", ""), 
            'port': 443
        }],
        http_compress = True,
        http_auth=(opensearch_account, opensearch_passwd),
        use_ssl = True,
        verify_certs = True,
        ssl_assert_hostname = False,
        ssl_show_warn = False,
    )

    if client.indices.exists(index_name):
        print('remove index: ', index_name)
        response = client.indices.delete(
            index=index_name
        )
        print('response(remove): ', response)    
    else:
        print('no index: ', index_name)

def delete_document_if_exist(vectorstore, metadata_key):
    try: 
        s3r = boto3.resource("s3")
        bucket = s3r.Bucket(s3_bucket)
        objs = list(bucket.objects.filter(Prefix=metadata_key))
        print('objs: ', objs)
        
        if(len(objs)>0):
            doc = s3r.Object(s3_bucket, metadata_key)
            meta = doc.get()['Body'].read().decode('utf-8')
            print('meta: ', meta)
            
            ids = json.loads(meta)['ids']
            print('ids: ', ids)
            
            result = vectorstore.delete(ids)
            print('result: ', result)        
        else:
            print('no meta file: ', metadata_key)
            
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        raise Exception ("Not able to create meta file")
    
def store_document_for_opensearch(bedrock_embedding, docs, key):
    index_name = 'idx-rag'
    
    vectorstore = OpenSearchVectorSearch(
            index_name=index_name,  
            is_aoss = False,
            #engine="faiss",  # default: nmslib
            embedding_function = bedrock_embedding,
            opensearch_url = opensearch_url,
            http_auth=(opensearch_account, opensearch_passwd),
        )
    
    objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
    print('objectName: ', objectName)    
    metadata_key = meta_prefix+objectName+'.metadata.json'
    print('meta file name: ', metadata_key)    
    delete_document_if_exist(vectorstore, metadata_key)
    
    try:
        response = vectorstore.add_documents(docs)
        print('response of adding documents: ', response)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                
        raise Exception ("Not able to request to LLM")

    print('uploaded into opensearch')
    
    return response
    
# load documents from s3 for pdf and txt
def load_document(file_type, s3_file_name):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, s3_prefix+'/'+s3_file_name)
    
    if file_type == 'pdf':
        Byte_contents = doc.get()['Body'].read()
        reader = PyPDF2.PdfReader(BytesIO(Byte_contents))
        
        texts = []
        for page in reader.pages:
            texts.append(page.extract_text())
        contents = '\n'.join(texts)
        
    elif file_type == 'pptx':
        Byte_contents = doc.get()['Body'].read()
            
        from pptx import Presentation
        prs = Presentation(BytesIO(Byte_contents))

        texts = []
        for i, slide in enumerate(prs.slides):
            text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = text + shape.text
            texts.append(text)
        contents = '\n'.join(texts)
        
    elif file_type == 'txt':        
        contents = doc.get()['Body'].read().decode('utf-8')

    elif file_type == 'docx':
        Byte_contents = doc.get()['Body'].read()
            
        import docx
        doc_contents =docx.Document(BytesIO(Byte_contents))

        texts = []
        for i, para in enumerate(doc_contents.paragraphs):
            if(para.text):
                texts.append(para.text)
                # print(f"{i}: {para.text}")        
        contents = '\n'.join(texts)
            
    # print('contents: ', contents)
    new_contents = str(contents).replace("\n"," ") 
    print('length: ', len(new_contents))

    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=100,
        separators=["\n\n", "\n", ".", " ", ""],
        length_function = len,
    ) 

    texts = text_splitter.split_text(new_contents) 
                
    return texts

# load csv documents from s3
def load_csv_document(path, doc_prefix, s3_file_name):
    s3r = boto3.resource("s3")
    doc = s3r.Object(s3_bucket, s3_prefix+'/'+s3_file_name)

    lines = doc.get()['Body'].read().decode('utf-8').split('\n')   # read csv per line
    print('lins: ', len(lines))
        
    columns = lines[0].split(',')  # get columns
    #columns = ["Category", "Information"]  
    #columns_to_metadata = ["type","Source"]
    print('columns: ', columns)
    
    docs = []
    n = 0
    for row in csv.DictReader(lines, delimiter=',',quotechar='"'):
        # print('row: ', row)
        #to_metadata = {col: row[col] for col in columns_to_metadata if col in row}
        values = {k: row[k] for k in columns if k in row}
        content = "\n".join(f"{k.strip()}: {v.strip()}" for k, v in values.items())
        doc = Document(
            page_content=content,
            metadata={
                'name': s3_file_name,
                'page': n+1,
                'uri': path+doc_prefix+parse.quote(s3_file_name)
            }
            #metadata=to_metadata
        )
        docs.append(doc)
        n = n+1
    print('docs[0]: ', docs[0])

    return docs

def get_summary(chat, docs):    
    text = ""
    for doc in docs:
        text = text + doc
    
    if isKorean(text)==True:
        system = (
            "다음의 <article> tag안의 문장을 요약해서 500자 이내로 설명하세오."
        )
    else: 
        system = (
            "Here is pieces of article, contained in <article> tags. Write a concise summary within 500 characters."
        )
    
    human = "<article>{text}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
    
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "text": text
            }
        )
        
        summary = result.content
        print('result of summarization: ', summary)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")
    
    return summary
    
def load_chat_history(userId, allowTime):
    dynamodb_client = boto3.client('dynamodb')

    response = dynamodb_client.query(
        TableName=callLogTableName,
        KeyConditionExpression='user_id = :userId AND request_time > :allowTime',
        ExpressionAttributeValues={
            ':userId': {'S': userId},
            ':allowTime': {'S': allowTime}
        }
    )
    # print('query result: ', response['Items'])

    for item in response['Items']:
        text = item['body']['S']
        msg = item['msg']['S']
        type = item['type']['S']

        if type == 'text':
            memory_chain.chat_memory.add_user_message(text)
            if len(msg) > MSG_LENGTH:
                memory_chain.chat_memory.add_ai_message(msg[:MSG_LENGTH])                          
            else:
                memory_chain.chat_memory.add_ai_message(msg) 
                
def getAllowTime():
    d = datetime.datetime.now() - datetime.timedelta(days = 2)
    timeStr = str(d)[0:19]
    print('allow time: ',timeStr)

    return timeStr

def isTyping(connectionId, requestId):    
    msg_proceeding = {
        'request_id': requestId,
        'msg': 'Proceeding...',
        'status': 'istyping'
    }
    #print('result: ', json.dumps(result))
    sendMessage(connectionId, msg_proceeding)

def readStreamMsg(connectionId, requestId, stream):
    msg = ""
    if stream:
        for event in stream:
            #print('event: ', event)
            msg = msg + event

            result = {
                'request_id': requestId,
                'msg': msg,
                'status': 'proceeding'
            }
            #print('result: ', json.dumps(result))
            sendMessage(connectionId, result)
    # print('msg: ', msg)
    return msg

def revise_question(connectionId, requestId, chat, query):    
    if isKorean(query)==True :      
        system = (
            ""
        )  
        human = """이전 대화를 참조하여, 다음의 <question>의 뜻을 명확히 하는 새로운 질문을 한국어로 생성하세요. 새로운 질문은 원래 질문의 중요한 단어를 반드시 포함합니다. 결과는 <result> tag를 붙여주세요.
        
        <question>            
        {question}
        </question>"""
        
    else: 
        system = (
            ""
        )
        human = """Rephrase the follow up <question> to be a standalone question. Put it in <result> tags.
        <question>            
        {question}
        </question>"""
            
    prompt = ChatPromptTemplate.from_messages([("system", system), MessagesPlaceholder(variable_name="history"), ("human", human)])
    print('prompt: ', prompt)
    
    history = memory_chain.load_memory_variables({})["chat_history"]
    print('memory_chain: ', history)
                
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "history": history,
                "question": query,
            }
        )
        generated_question = result.content
        
        revised_question = generated_question[generated_question.find('<result>')+8:len(generated_question)-9] # remove <result> tag                   
        print('revised_question: ', revised_question)
        
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
            
        sendErrorMessage(connectionId, requestId, err_msg)    
        raise Exception ("Not able to request to LLM")
            
    return revised_question    
    # return revised_question.replace("\n"," ")

def query_using_RAG_context(connectionId, requestId, chat, context, revised_question):    
    if isKorean(revised_question)==True:
        system = (
            """다음의 <context> tag안의 참고자료를 이용하여 상황에 맞는 구체적인 세부 정보를 충분히 제공합니다. Assistant의 이름은 서연이고, 모르는 질문을 받으면 솔직히 모른다고 말합니다.
            
            <context>
            {context}
            </context>"""
        )
    else: 
        system = (
            """Here is pieces of context, contained in <context> tags. Provide a concise answer to the question at the end. If you don't know the answer, just say that you don't know, don't try to make up an answer.
            
            <context>
            {context}
            </context>"""
        )
    
    human = "{input}"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
                   
    chain = prompt | chat
    
    try: 
        isTyping(connectionId, requestId)  
        stream = chain.invoke(
            {
                "context": context,
                "input": revised_question,
            }
        )
        msg = readStreamMsg(connectionId, requestId, stream.content)    
        print('msg: ', msg)
        
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
            
        sendErrorMessage(connectionId, requestId, err_msg)    
        raise Exception ("Not able to request to LLM")

    return msg

def priority_search(query, relevant_docs, bedrock_embedding):
    excerpts = []
    for i, doc in enumerate(relevant_docs):
        # print('doc: ', doc)
        if doc['metadata']['translated_excerpt']:
            content = doc['metadata']['translated_excerpt']
        else:
            content = doc['metadata']['excerpt']
        
        excerpts.append(
            Document(
                page_content=content,
                metadata={
                    'name': doc['metadata']['title'],
                    'order':i,
                }
            )
        )  
    # print('excerpts: ', excerpts)

    embeddings = bedrock_embedding
    vectorstore_confidence = FAISS.from_documents(
        excerpts,  # documents
        embeddings  # embeddings
    )            
    rel_documents = vectorstore_confidence.similarity_search_with_score(
        query=query,
        k=top_k
    )

    docs = []
    for i, document in enumerate(rel_documents):
        print(f'## Document(priority_search) {i+1}: {document}')

        order = document[0].metadata['order']
        name = document[0].metadata['name']
        assessed_score = document[1]
        print(f"{order} {name}: {assessed_score}")

        relevant_docs[order]['assessed_score'] = int(assessed_score)

        if assessed_score < 200:
            docs.append(relevant_docs[order])    
    # print('selected docs: ', docs)

    return docs

def get_reference(docs, path, doc_prefix):
    reference = "\n\nFrom\n"
    for i, doc in enumerate(docs):
        if doc['metadata']['translated_excerpt']:
            excerpt = str(doc['metadata']['excerpt']+'  [번역]'+doc['metadata']['translated_excerpt']).replace('"'," ") 
        else:
            excerpt = str(doc['metadata']['excerpt']).replace('"'," ")

        if doc['rag_type'] == 'opensearch':
            print(f'## Document(get_reference) {i+1}: {doc}')
                
            page = ""
            if "document_attributes" in doc['metadata']:
                if "_excerpt_page_number" in doc['metadata']['document_attributes']:
                    page = doc['metadata']['document_attributes']['_excerpt_page_number']
            uri = doc['metadata']['source']
            name = doc['metadata']['title']

            #print('opensearch page: ', page)

            if page:                
                reference = reference + f"{i+1}. {page}page in <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']})\n"
            else:
                reference = reference + f"{i+1}. <a href={uri} target=_blank>{name}</a>, {doc['rag_type']} ({doc['assessed_score']}), <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
                    
        elif doc['rag_type'] == 'search': # google search
            print(f'## Document(get_reference) {i+1}: {doc}')
                
            uri = doc['metadata']['source']
            name = doc['metadata']['title']
            reference = reference + f"{i+1}. <a href={uri} target=_blank>{name}</a>, {doc['rag_type']}, <a href=\"#\" onClick=\"alert(`{excerpt}`)\">관련문서</a>\n"
        
    return reference

def retrieve_from_vectorstore(query, top_k, rag_type):
    print(f"query: {query} ({rag_type})")

    relevant_docs = []
    if rag_type == 'opensearch':
        relevant_documents = vectorstore_opensearch.similarity_search_with_score(
            query = query,
            k = top_k,
        )
        #print('(opensearch score) relevant_documents: ', relevant_documents)

        for i, document in enumerate(relevant_documents):
            print(f'## Document(retrieve_from_vectorstore) {i+1}: {document}')

            name = document[0].metadata['name']
            print('metadata: ', document[0].metadata)

            page = ""
            if "page" in document[0].metadata:
                page = document[0].metadata['page']
            uri = ""
            if "uri" in document[0].metadata:
                uri = document[0].metadata['uri']

            excerpt = document[0].page_content
            confidence = str(document[1])
            assessed_score = str(document[1])

            if page:
                print('page: ', page)
                doc_info = {
                    "rag_type": rag_type,
                    "confidence": confidence,
                    "metadata": {
                        "source": uri,
                        "title": name,
                        "excerpt": excerpt,
                        "translated_excerpt": "",
                        "document_attributes": {
                            "_excerpt_page_number": page
                        }
                    },
                    "assessed_score": assessed_score,
                }
            else:
                doc_info = {
                    "rag_type": rag_type,
                    "confidence": confidence,
                    "metadata": {
                        "source": uri,
                        "title": name,
                        "excerpt": excerpt,
                        "translated_excerpt": ""
                    },
                    "assessed_score": assessed_score,
                }
            relevant_docs.append(doc_info)

    return relevant_docs

def translate_process_from_relevent_doc(conn, chat, doc):
    translated_excerpt = translate_text(chat, doc['metadata']['excerpt'])

    # doc['metadata']['excerpt'] = translated_excerpt
    doc['metadata']['translated_excerpt'] = translated_excerpt

    conn.send(doc)
    conn.close()

def translate_relevant_documents_using_parallel_processing(docs):
    selected_LLM = 0
    relevant_docs = []    
    processes = []
    parent_connections = []
    for doc in docs:
        parent_conn, child_conn = Pipe()
        parent_connections.append(parent_conn)
            
        chat = get_chat(profile_of_LLMs, selected_LLM)
        process = Process(target=translate_process_from_relevent_doc, args=(child_conn, chat, doc))            
        processes.append(process)

        selected_LLM = selected_LLM + 1
        if selected_LLM == len(profile_of_LLMs):
            selected_LLM = 0

    for process in processes:
        process.start()
            
    for parent_conn in parent_connections:
        doc = parent_conn.recv()
        relevant_docs.append(doc)    

    for process in processes:
        process.join()
    
    #print('relevant_docs: ', relevant_docs)
    return relevant_docs

def get_answer_using_RAG(chat, text, conv_type, connectionId, requestId, bedrock_embedding):
    global time_for_revise, time_for_rag, time_for_inference, time_for_priority_search, number_of_relevant_docs  # for debug
    time_for_revise = time_for_rag = time_for_inference = time_for_priority_search = number_of_relevant_docs = 0
    
    reference = ""
    start_time_for_revise = time.time()

    # revise question
    revised_question = revise_question(connectionId, requestId, chat, text)     
    print('revised_question: ', revised_question)     

    end_time_for_revise = time.time()
    time_for_revise = end_time_for_revise - start_time_for_revise
    print('processing time for revised question: ', time_for_revise)

    relevant_docs = [] 
    print('start RAG for revised question')

    rag_type = 'opensearch'
    relevant_docs = retrieve_from_vectorstore(query=revised_question, top_k=top_k, rag_type=rag_type)
    print(f'relevant_docs ({rag_type}): '+json.dumps(relevant_docs))

    if allowDualSearching=='true' and isKorean(text)==True:
        print('start RAG for translated revised question')
        translated_revised_question = translate_text(chat, revised_question)
        print('translated_revised_question: ', translated_revised_question)

        relevant_docs_using_translated_question = retrieve_from_vectorstore(query=translated_revised_question, top_k=4, rag_type=rag_type)
                    
        docs_translation_required = []
        if len(relevant_docs_using_translated_question)>=1:
            for i, doc in enumerate(relevant_docs_using_translated_question):
                if isKorean(doc)==False:
                    docs_translation_required.append(doc)
                else:
                    print(f"original {i}: {doc}")
                    relevant_docs.append(doc)
                                           
            translated_docs = translate_relevant_documents_using_parallel_processing(docs_translation_required)
            for i, doc in enumerate(translated_docs):
                print(f"#### {i} (ENG): {doc['metadata']['excerpt']}")
                print(f"#### {i} (KOR): {doc['metadata']['translated_excerpt']}")
                relevant_docs.append(doc)

    end_time_for_rag = time.time()
    time_for_rag = end_time_for_rag - end_time_for_revise
    print('processing time for RAG: ', time_for_rag)

    selected_relevant_docs = []
    if len(relevant_docs)>=1:
        selected_relevant_docs = priority_search(revised_question, relevant_docs, bedrock_embedding)
        print('selected_relevant_docs: ', json.dumps(selected_relevant_docs))

    if len(selected_relevant_docs)==0 and allowInternetSearching=='true':
        print('No relevant document! So use google api')            
        api_key = google_api_key
        cse_id = google_cse_id 
            
        relevant_docs = []
        try: 
            service = build("customsearch", "v1", developerKey=api_key)
            result = service.cse().list(q=revised_question, cx=cse_id).execute()
            print('google search result: ', result)

            if "items" in result:
                for item in result['items']:
                    api_type = "google api"
                    excerpt = item['snippet']
                    uri = item['link']
                    title = item['title']
                    confidence = ""
                    assessed_score = ""
                        
                    doc_info = {
                        "rag_type": 'search',
                        "api_type": api_type,
                        "confidence": confidence,
                        "metadata": {
                            "source": uri,
                            "title": title,
                            "excerpt": excerpt,
                            "translated_excerpt": "",
                        },
                        "assessed_score": assessed_score,
                    }
                    relevant_docs.append(doc_info)                
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)       

            sendErrorMessage(connectionId, requestId, err_msg)    
            raise Exception ("Not able to search using google api")   
            
        if len(relevant_docs)>=1:
            selected_relevant_docs = priority_search(revised_question, relevant_docs, bedrock_embedding)
            print('selected_relevant_docs: ', json.dumps(selected_relevant_docs))
        print('selected_relevant_docs (google): ', selected_relevant_docs)

    end_time_for_priority_search = time.time() 
    time_for_priority_search = end_time_for_priority_search - end_time_for_rag
    print('processing time for priority search: ', time_for_priority_search)
    number_of_relevant_docs = len(selected_relevant_docs)

    relevant_context = ""
    for document in selected_relevant_docs:
        if document['metadata']['translated_excerpt']:
            content = document['metadata']['translated_excerpt']
        else:
            content = document['metadata']['excerpt']
        
        relevant_context = relevant_context + content + "\n\n"
    print('relevant_context: ', relevant_context)

    # query using RAG context
    msg = query_using_RAG_context(connectionId, requestId, chat, relevant_context, revised_question)

    reference = ""
    if len(selected_relevant_docs)>=1 and enableReference=='true':
        reference = get_reference(selected_relevant_docs, path, doc_prefix)  

    end_time_for_inference = time.time()
    time_for_inference = end_time_for_inference - end_time_for_priority_search
    print('processing time for inference: ', time_for_inference)

    return msg, reference

def translate_text(chat, text):
    system = (
        "You are a helpful assistant that translates {input_language} to {output_language} in <article> tags. Put it in <result> tags."
    )
    human = "<article>{text}</article>"
    
    prompt = ChatPromptTemplate.from_messages([("system", system), ("human", human)])
    print('prompt: ', prompt)
    
    if isKorean(text)==False :
        input_language = "English"
        output_language = "Korean"
    else:
        input_language = "Korean"
        output_language = "English"
                        
    chain = prompt | chat    
    try: 
        result = chain.invoke(
            {
                "input_language": input_language,
                "output_language": output_language,
                "text": text,
            }
        )
        
        msg = result.content
        print('translated text: ', msg)
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)                    
        raise Exception ("Not able to request to LLM")

    return msg[msg.find('<result>')+8:len(msg)-9] # remove <result> tag

def create_metadata(bucket, key, meta_prefix, s3_prefix, uri, category, documentId, ids):
    title = key
    timestamp = int(time.time())

    metadata = {
        "Attributes": {
            "_category": category,
            "_source_uri": uri,
            "_version": str(timestamp),
            "_language_code": "ko"
        },
        "Title": title,
        "DocumentId": documentId,      
        "ids": ids  
    }
    print('metadata: ', metadata)
    
    #objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)]).upper()
    objectName = (key[key.find(s3_prefix)+len(s3_prefix)+1:len(key)])
    print('objectName: ', objectName)

    client = boto3.client('s3')
    try: 
        client.put_object(
            Body=json.dumps(metadata), 
            Bucket=bucket, 
            Key=meta_prefix+objectName+'.metadata.json' 
        )
    except Exception:
        err_msg = traceback.format_exc()
        print('error message: ', err_msg)        
        raise Exception ("Not able to create meta file")

def getResponse(connectionId, jsonBody):
    userId  = jsonBody['user_id']
    # print('userId: ', userId)
    requestId  = jsonBody['request_id']
    # print('requestId: ', requestId)
    requestTime  = jsonBody['request_time']
    # print('requestTime: ', requestTime)
    type  = jsonBody['type']
    # print('type: ', type)
    body = jsonBody['body']
    # print('body: ', body)
    conv_type = jsonBody['conv_type']  # conversation type
    print('Conversation Type: ', conv_type)

    global vectorstore_opensearch, enableReference
    global map_chain, memory_chain, isReady, selected_LLM, allowDualSearching, allowInternetSearching

    # Multi-LLM
    profile = profile_of_LLMs[selected_LLM]
    bedrock_region =  profile['bedrock_region']
    modelId = profile['model_id']
    print(f'selected_LLM: {selected_LLM}, bedrock_region: {bedrock_region}, modelId: {modelId}')
    # print('profile: ', profile)
    
    chat = get_chat(profile_of_LLMs, selected_LLM)    
    bedrock_embedding = get_embedding(profile_of_LLMs, selected_LLM)

    # allocate memory
    if userId in map_chain:  
        print('memory exist. reuse it!')        
        memory_chain = map_chain[userId]
        
    else: 
        print('memory does not exist. create new one!')
        memory_chain = ConversationBufferWindowMemory(memory_key="chat_history", output_key='answer', return_messages=True, k=10)
        map_chain[userId] = memory_chain
        
        allowTime = getAllowTime()
        load_chat_history(userId, allowTime)

    # rag sources
    if conv_type == 'qa':
        vectorstore_opensearch = OpenSearchVectorSearch(
            index_name = "idx-rag*", # all
            is_aoss = False,
            ef_search = 1024, # 512(default)
            m=48,
            embedding_function = bedrock_embedding,
            opensearch_url=opensearch_url,
            http_auth=(opensearch_account, opensearch_passwd), # http_auth=awsauth,
        )

    start = int(time.time())    

    msg = ""
    reference = ""
    if type == 'text' and body[:11] == 'list models':
        bedrock_client = boto3.client(
            service_name='bedrock',
            region_name=bedrock_region,
        )
        modelInfo = bedrock_client.list_foundation_models()    
        print('models: ', modelInfo)

        msg = f"The list of models: \n"
        lists = modelInfo['modelSummaries']
        
        for model in lists:
            msg += f"{model['modelId']}\n"
        
        msg += f"current model: {modelId}"
        print('model lists: ', msg)          

        sendResultMessage(connectionId, requestId, msg)  
    else:             
        if type == 'text':
            text = body
            print('query: ', text)

            querySize = len(text)
            textCount = len(text.split())
            print(f"query size: {querySize}, words: {textCount}")

            if text == 'enableReference':
                enableReference = 'true'
                msg  = "Referece is enabled"
            elif text == 'disableReference':
                enableReference = 'false'
                msg  = "Reference is disabled"
            elif text == 'enableDualSearching':
                allowDualSearching = 'true'
                msg  = "Translated question is enabled"
            elif text == 'disableDualSearching':
                allowDualSearching = 'false'
                msg  = "Translated question is disabled"
            elif text == 'enableInternetSearching':
                allowInternetSearching = 'true'
                msg  = "Internet Search is enabled"
            elif text == 'disableInternetSearching':
                allowInternetSearching = 'false'
                msg  = "Internet Search is disabled"

            elif text == 'clearMemory':
                memory_chain.clear()
                map_chain[userId] = memory_chain
                    
                print('initiate the chat memory!')
                msg  = "The chat memory was intialized in this session."
            else:          
                if conv_type == 'normal':
                    msg = general_conversation(connectionId, requestId, chat, text)   
                if conv_type == 'qa':   # RAG
                    msg, reference = get_answer_using_RAG(chat, text, conv_type, connectionId, requestId, bedrock_embedding)     
            
            memory_chain.chat_memory.add_user_message(text)  # append new diaglog
            memory_chain.chat_memory.add_ai_message(msg)
                
        elif type == 'document':
            isTyping(connectionId, requestId)
            
            object = body
            file_type = object[object.rfind('.')+1:len(object)]            
            print('file_type: ', file_type)

            ids = []
            if file_type == 'csv':
                docs = load_csv_document(path, doc_prefix, object)
                contexts = []
                for doc in docs:
                    contexts.append(doc.page_content)
                print('contexts: ', contexts)

                msg = get_summary(chat, contexts)

            elif file_type == 'pdf' or file_type == 'txt' or file_type == 'pptx' or file_type == 'docx':
                texts = load_document(file_type, object)

                docs = []
                for i in range(len(texts)):
                    docs.append(
                        Document(
                            page_content=texts[i],
                            metadata={
                                'name': object,
                                # 'page':i+1,
                                'uri': path+doc_prefix+parse.quote(object)
                            }
                        )
                    )
                print('docs[0]: ', docs[0])    
                print('docs size: ', len(docs))

                contexts = []
                for doc in docs:
                    contexts.append(doc.page_content)
                print('contexts: ', contexts)

                msg = get_summary(chat, contexts)
            else:
                msg = "uploaded file: "+object
                                
            if conv_type == 'qa':
                start_time = time.time()
                
                category = "upload"
                key = object
                documentId = category + "-" + key
                documentId = documentId.replace(' ', '_') # remove spaces
                documentId = documentId.replace(',', '_') # remove commas
                documentId = documentId.lower() # change to lowercase
                print('documentId: ', documentId)

                if file_type == 'pdf' or file_type == 'txt' or file_type == 'csv' or file_type == 'pptx' or file_type == 'docx':
                    ids = store_document_for_opensearch(bedrock_embedding, docs, key)
                
                print('processing time: ', str(time.time() - start_time))
            
                create_metadata(bucket=s3_bucket, key=key, meta_prefix=meta_prefix, s3_prefix=s3_prefix, uri=path+parse.quote(key), category=category, documentId=documentId, ids=ids)
                        
        sendResultMessage(connectionId, requestId, msg+reference)
        # print('msg+reference: ', msg+reference)

        elapsed_time = time.time() - start
        print("total run time(sec): ", elapsed_time)

        if isKorean(msg)==False and conv_type=='qa':
            translated_msg = translate_text(chat, msg)
            print('translated_msg: ', translated_msg)

            msg = msg+'\n[한국어]\n'+translated_msg
            sendResultMessage(connectionId, requestId, msg+reference)

        item = {    # save dialog
            'user_id': {'S':userId},
            'request_id': {'S':requestId},
            'request_time': {'S':requestTime},
            'type': {'S':type},
            'body': {'S':body},
            'msg': {'S':msg+reference}
        }
        client = boto3.client('dynamodb')
        try:
            resp =  client.put_item(TableName=callLogTableName, Item=item)
        except Exception:
            err_msg = traceback.format_exc()
            print('error message: ', err_msg)
            raise Exception ("Not able to write into dynamodb")        
        #print('resp, ', resp)

    if selected_LLM >= len(profile_of_LLMs)-1:
        selected_LLM = 0
    else:
        selected_LLM = selected_LLM + 1

    return msg, reference

def lambda_handler(event, context):
    # print('event: ', event)
    
    msg = ""
    if event['requestContext']: 
        connectionId = event['requestContext']['connectionId']        
        routeKey = event['requestContext']['routeKey']
        
        if routeKey == '$connect':
            print('connected!')
        elif routeKey == '$disconnect':
            print('disconnected!')
        else:
            body = event.get("body", "")
            #print("data[0:8]: ", body[0:8])
            if body[0:8] == "__ping__":
                # print("keep alive!")                
                sendMessage(connectionId, "__pong__")
            else:
                print('connectionId: ', connectionId)
                print('routeKey: ', routeKey)
        
                jsonBody = json.loads(body)
                print('request body: ', json.dumps(jsonBody))

                requestId  = jsonBody['request_id']
                try:
                    msg, reference = getResponse(connectionId, jsonBody)

                    print('msg+reference: ', msg+reference)
                                        
                except Exception:
                    err_msg = traceback.format_exc()
                    print('err_msg: ', err_msg)

                    sendErrorMessage(connectionId, requestId, err_msg)    
                    raise Exception ("Not able to send a message")

    return {
        'statusCode': 200
    }
